"""PPTX parser that extracts slide text and speaker notes."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Dict, Iterable, Iterator, Optional

from pptx import Presentation  # type: ignore

LOGGER = logging.getLogger(__name__)


def _slide_text(slide) -> Iterator[str]:  # type: ignore[no-untyped-def]
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = (shape.text or "").strip()
        if text:
            yield text


def _notes_text(slide) -> Optional[str]:  # type: ignore[no-untyped-def]
    if not slide.has_notes_slide:  # type: ignore[attr-defined]
        return None
    notes_slide = slide.notes_slide  # type: ignore[attr-defined]
    if not hasattr(notes_slide, "notes_text_frame"):
        return None
    notes = notes_slide.notes_text_frame.text or ""
    cleaned = notes.strip()
    return cleaned or None


def parse_pptx(path: str, *, metadata: Optional[Dict[str, str]] = None) -> Iterable[Dict[str, object]]:
    """Yield raw element dictionaries produced from a PPTX file."""
    meta = metadata or {}
    pptx_path = Path(path)
    doc_title = meta.get("doc_title", pptx_path.stem.replace("_", " "))

    presentation = Presentation(str(pptx_path))
    for slide_index, slide in enumerate(presentation.slides, start=1):
        source_uri = f"file://{pptx_path.resolve()}#slide={slide_index}"
        for block_index, text in enumerate(_slide_text(slide), start=1):
            yield {
                "doc_path": str(pptx_path.resolve()),
                "doc_title": doc_title,
                "doc_type": "pptx",
                "element_type": "paragraph",
                "order": block_index,
                "content_text": text,
                "source_uri": f"{source_uri}&shape={block_index}",
                "metadata": meta,
            }

        notes = _notes_text(slide)
        if notes:
            yield {
                "doc_path": str(pptx_path.resolve()),
                "doc_title": doc_title,
                "doc_type": "pptx",
                "element_type": "note",
                "order": 0,
                "content_text": notes,
                "source_uri": f"{source_uri}&notes=1",
                "metadata": meta,
            }
